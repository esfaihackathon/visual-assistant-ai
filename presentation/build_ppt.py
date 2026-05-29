"""
Saral — Accessible Banking Assistant
Full PowerPoint presentation generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import io, math

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x08, 0x13, 0x20)
NAVY_CARD  = RGBColor(0x0E, 0x20, 0x35)
NAVY_LIGHT = RGBColor(0x14, 0x2A, 0x45)
BLUE       = RGBColor(0x4D, 0x9E, 0xFF)
TEAL       = RGBColor(0x00, 0xBF, 0xA5)
AMBER      = RGBColor(0xFF, 0xC1, 0x07)
ORANGE     = RGBColor(0xFF, 0x6B, 0x35)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTED      = RGBColor(0xB2, 0xCC, 0xE8)
DARK_TEXT  = RGBColor(0x08, 0x13, 0x20)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank

# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    return txb

def add_multiline(slide, lines, x, y, w, h, font_size=16, color=MUTED,
                  align=PP_ALIGN.LEFT, spacing_after=Pt(6), font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (txt, bold, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = spacing_after
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = clr if clr else color
        run.font.name = font_name
    return txb

def navy_bg(slide):
    add_rect(slide, 0, 0, W, H, fill=NAVY)

def section_label(slide, text, x=Inches(0.6), y=Inches(0.35)):
    add_textbox(slide, text, x, y, Inches(10), Inches(0.4),
                font_size=11, bold=True, color=AMBER,
                font_name="Calibri")

def slide_title(slide, text, x=Inches(0.6), y=Inches(0.7), w=Inches(12), h=Inches(1)):
    add_textbox(slide, text, x, y, w, h,
                font_size=34, bold=True, color=WHITE, font_name="Calibri")

def accent_bar(slide, y=Inches(1.62)):
    add_rect(slide, Inches(0.6), y, Inches(1.4), Inches(0.055), fill=BLUE)

def card_rect(slide, x, y, w, h, accent_color=None):
    r = add_rect(slide, x, y, w, h, fill=NAVY_CARD, line=NAVY_LIGHT, line_w=Pt(1))
    if accent_color:
        add_rect(slide, x, y, Inches(0.07), h, fill=accent_color)
    return r

def stat_block(slide, x, y, w, h, number, label, num_color=BLUE):
    card_rect(slide, x, y, w, h)
    add_textbox(slide, number, x, y + Inches(0.22), w, Inches(0.7),
                font_size=36, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x, y + Inches(0.9), w, Inches(0.6),
                font_size=12, color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)

# Gradient accent panel left
add_rect(sl, 0, 0, Inches(4.5), H, fill=RGBColor(0x0A, 0x1A, 0x32))
add_rect(sl, 0, 0, Inches(0.12), H, fill=BLUE)

# App icon circle
add_rect(sl, Inches(1.5), Inches(1.1), Inches(1.5), Inches(1.5), fill=RGBColor(0x10, 0x2A, 0x4A))
add_textbox(sl, "🏦", Inches(1.5), Inches(1.15), Inches(1.5), Inches(1.4),
            font_size=54, align=PP_ALIGN.CENTER)

add_textbox(sl, "SARAL", Inches(0.45), Inches(2.85), Inches(3.6), Inches(1.1),
            font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, "सरल", Inches(0.45), Inches(3.7), Inches(3.6), Inches(0.5),
            font_size=20, color=AMBER, align=PP_ALIGN.CENTER, italic=True)
add_rect(sl, Inches(0.6), Inches(4.2), Inches(3.1), Inches(0.05), fill=BLUE)
add_textbox(sl, "Voice-First Accessible Banking", Inches(0.3), Inches(4.35), Inches(3.9), Inches(0.5),
            font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Right panel — key points
rx = Inches(5.0)
add_textbox(sl, "ESFAI Hackathon  •  2026", rx, Inches(1.0), Inches(7.8), Inches(0.4),
            font_size=12, bold=True, color=AMBER)
add_textbox(sl, "Banking for everyone —\nregardless of sight, age, or literacy.",
            rx, Inches(1.55), Inches(7.8), Inches(1.2),
            font_size=26, bold=True, color=WHITE, wrap=True)
add_rect(sl, rx, Inches(2.85), Inches(1.5), Inches(0.06), fill=TEAL)

bullets = [
    ("🎙️  100% voice-operable", TEAL),
    ("♿  WCAG AAA compliant",   BLUE),
    ("👆  Auto fingerprint auth", AMBER),
    ("📍  Fixed mic — always reachable", ORANGE),
    ("🌐  Android + Web",        MUTED),
]
for i, (txt, clr) in enumerate(bullets):
    add_textbox(sl, txt, rx, Inches(3.15) + Inches(0.54)*i,
                Inches(7.5), Inches(0.48), font_size=16, color=clr, bold=False)

# Bottom bar
add_rect(sl, 0, Inches(7.0), W, Inches(0.5), fill=RGBColor(0x05, 0x0D, 0x18))
add_textbox(sl, "v8.0  •  439 Tests Passing  •  0 Failures  •  Android + Web",
            Inches(0.3), Inches(7.05), W - Inches(0.6), Inches(0.38),
            font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "THE PROBLEM")
slide_title(sl, "Banking apps leave hundreds of millions behind")
accent_bar(sl)

# Problem cards — 2×2
problems = [
    ("👁️", "Visual Barriers",
     "~72M Indians live with vision impairment or blindness (IAPB Vision Atlas 2023). Small fonts, low contrast, colour-coded status — completely inaccessible.",
     BLUE),
    ("🤲", "Motor Barriers",
     "~145M senior citizens today (347M projected by 2050) struggle with precise taps on tiny buttons. Auto-biometric removes the biggest friction point.",
     TEAL),
    ("📖", "Literacy Barriers",
     "~372M Indians have limited literacy (Census 2011, 74% literacy rate on 1.43B population) — cannot navigate menu-driven interfaces.",
     AMBER),
    ("🎨", "Colour-Blindness",
     "~300M people globally are colour-blind; ~60M in India (~8% of males). Red/green status indicators exclude them entirely.",
     ORANGE),
]
positions = [
    (Inches(0.5), Inches(1.85)),
    (Inches(6.65), Inches(1.85)),
    (Inches(0.5), Inches(4.35)),
    (Inches(6.65), Inches(4.35)),
]
for (x, y), (icon, title, body, clr) in zip(positions, problems):
    card_rect(sl, x, y, Inches(6.0), Inches(2.3), accent_color=clr)
    add_textbox(sl, icon, x + Inches(0.25), y + Inches(0.15), Inches(0.8), Inches(0.7), font_size=28)
    add_textbox(sl, title, x + Inches(1.1), y + Inches(0.15), Inches(4.7), Inches(0.45),
                font_size=17, bold=True, color=WHITE)
    add_textbox(sl, body, x + Inches(0.22), y + Inches(0.72), Inches(5.6), Inches(1.4),
                font_size=13, color=MUTED, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHO IS IMPACTED (with bar chart)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "PEOPLE IMPACTED")
slide_title(sl, "Scale of the accessibility gap in India")
accent_bar(sl)

# Bar chart
cd = ChartData()
cd.categories = ["Visually\nImpaired", "Colour-\nBlind", "Senior\nCitizens (60+)", "Low-\nLiteracy"]
cd.add_series("Population (Millions)", (72, 60, 145, 372))

chart = sl.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.45), Inches(1.85), Inches(6.0), Inches(4.8), cd
).chart

chart.has_legend = False
chart.has_title  = False
plot = chart.plots[0]
plot.gap_width = 60

# Style series bars
series = plot.series[0]
fills = [BLUE, ORANGE, TEAL, AMBER]
for i, pt in enumerate(series.points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = fills[i]

chart.category_axis.tick_labels.font.size = Pt(12)
chart.category_axis.tick_labels.font.color.rgb = MUTED
chart.value_axis.tick_labels.font.size = Pt(11)
chart.value_axis.tick_labels.font.color.rgb = MUTED

# Stat panels on the right
stats = [
    ("~72M",  "Visually impaired\nIndians (IAPB 2023)",   BLUE),
    ("~60M",  "Colour-blind\nin India (~8% males)",      ORANGE),
    ("145M+", "Senior citizens 60+\n(347M by 2050, UN)", TEAL),
    ("~372M", "Limited literacy\n(Census 2011, 26%)",    AMBER),
]
for i, (num, lbl, clr) in enumerate(stats):
    sy = Inches(1.82) + Inches(1.18) * i
    stat_block(sl, Inches(6.65), sy, Inches(2.85), Inches(1.08), num, lbl, clr)

# Note
add_textbox(sl, "Sources: IAPB Vision Atlas 2023 · Census of India 2011 · UN World Population Ageing 2023 · Colour Blind Awareness (colourblindawareness.org)  |  See References slide",
            Inches(0.45), Inches(7.0), Inches(8), Inches(0.38),
            font_size=10, color=MUTED, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — OUR SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "OUR SOLUTION")
slide_title(sl, "Saral — Voice-First Banking Assistant")
accent_bar(sl)

add_textbox(sl, "Every banking action is completable by voice alone. The visual UI is a high-contrast fallback — not the primary interface.",
            Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.65),
            font_size=15, color=MUTED, wrap=True)

pillars = [
    ("🎙️", "Voice is Primary",
     "Balance • Transfer • Transactions • Cheque Book — all operable by speaking alone.",
     BLUE),
    ("👆", "Touch is Secondary",
     "88dp minimum tap targets. Large beneficiary cards. Mic button always fixed at the bottom.",
     TEAL),
    ("♿", "WCAG AAA Compliant",
     "18.7:1 contrast (WCAG AAA). Colour-blind safe palette. Full ARIA labels. Screen reader compatible.",
     AMBER),
    ("🔒", "Auto Biometric",
     "Fingerprint scanner activates automatically after transfer confirmation — no button tap needed.",
     ORANGE),
]
for i, (icon, title, body, clr) in enumerate(pillars):
    x = Inches(0.4) + Inches(3.2) * i
    card_rect(sl, x, Inches(2.55), Inches(3.02), Inches(3.5), accent_color=clr)
    add_textbox(sl, icon, x + Inches(0.15), Inches(2.65), Inches(0.8), Inches(0.7), font_size=30)
    add_textbox(sl, title, x + Inches(0.18), Inches(3.38), Inches(2.75), Inches(0.5),
                font_size=15, bold=True, color=WHITE)
    add_textbox(sl, body, x + Inches(0.18), Inches(3.85), Inches(2.78), Inches(2.0),
                font_size=12, color=MUTED, wrap=True)

# Comparison row
add_rect(sl, Inches(0.4), Inches(6.22), Inches(12.4), Inches(0.06), fill=NAVY_LIGHT)
compare = [
    ("Standard Banking App", "Requires reading, precise tapping, colour interpretation", ORANGE),
    ("Saral", "Operates entirely by voice — zero reading required", TEAL),
]
for i, (app, desc, clr) in enumerate(compare):
    x = Inches(0.5) + Inches(6.2) * i
    add_textbox(sl, app,  x, Inches(6.38), Inches(6.0), Inches(0.38), font_size=13, bold=True, color=clr)
    add_textbox(sl, desc, x, Inches(6.72), Inches(6.0), Inches(0.6),  font_size=12, color=MUTED, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — USER JOURNEY (Transfer Flow)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "USER EXPERIENCE")
slide_title(sl, "Complete a money transfer — entirely by voice")
accent_bar(sl)

steps = [
    ("🏠", "Home",      '"Transfer money"'),
    ("👥", "Select",    '"Rahul"  or tap name'),
    ("💰", "Amount",    '"500 rupees"'),
    ("✅", "Confirm",   '"Yes"'),
    ("👆", "Biometric", "Auto-scans finger"),
    ("🎉", "Complete",  "Balance read aloud"),
]
n   = len(steps)
sw  = Inches(1.72)
gap = Inches(0.26)
sx0 = Inches(0.42)
sy  = Inches(2.1)

colors_step = [BLUE, TEAL, AMBER, ORANGE, RGBColor(0xAB,0x47,0xBC), TEAL]

for i, (icon, title, caption) in enumerate(steps):
    x = sx0 + (sw + gap) * i
    clr = colors_step[i]
    # Card
    card_rect(sl, x, sy, sw, Inches(3.1))
    add_rect(sl, x, sy, sw, Inches(0.08), fill=clr)
    # Step number
    add_rect(sl, x + Inches(0.06), sy + Inches(0.15), Inches(0.38), Inches(0.38),
             fill=clr)
    add_textbox(sl, str(i+1), x + Inches(0.06), sy + Inches(0.14),
                Inches(0.38), Inches(0.38), font_size=13, bold=True,
                color=DARK_TEXT, align=PP_ALIGN.CENTER)
    add_textbox(sl, icon,   x, sy + Inches(0.6),  sw, Inches(0.7), font_size=32, align=PP_ALIGN.CENTER)
    add_textbox(sl, title,  x, sy + Inches(1.35), sw, Inches(0.45), font_size=14, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, caption, x, sy + Inches(1.82), sw, Inches(1.1), font_size=11,
                color=MUTED, align=PP_ALIGN.CENTER, wrap=True, italic=True)
    # Arrow
    if i < n - 1:
        ax = x + sw + Inches(0.02)
        add_textbox(sl, "→", ax, sy + Inches(1.25), gap + Inches(0.24), Inches(0.55),
                    font_size=22, color=BLUE, align=PP_ALIGN.CENTER)

# Key insight box
add_rect(sl, Inches(0.4), Inches(5.45), Inches(12.4), Inches(1.72), fill=RGBColor(0x0A,0x20,0x38))
add_rect(sl, Inches(0.4), Inches(5.45), Inches(0.08), Inches(1.72), fill=TEAL)
add_textbox(sl, "🔑  Key insight:",
            Inches(0.65), Inches(5.58), Inches(11.8), Inches(0.38),
            font_size=14, bold=True, color=TEAL)
add_textbox(sl,
    "At no step does the user need to read the screen, type, or find a small button. "
    "The fingerprint scanner activates automatically — the user simply rests their finger on the sensor.",
    Inches(0.65), Inches(5.95), Inches(12.0), Inches(1.1),
    font_size=14, color=WHITE, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FEATURES (grid)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "FEATURES")
slide_title(sl, "Everything a user needs — by voice")
accent_bar(sl)

features = [
    ("💳", "Check Balance",         "Speaks account name, last 4 digits, and available balance. Asks what to do next.",              BLUE),
    ("💸", "Transfer Money",        "End-to-end by voice: select beneficiary → amount → confirm → auto-biometric → complete.",        TEAL),
    ("📋", "Transaction History",   "Last 1/5/10 transactions read aloud. Follow-up asks 'more or main menu?'",                      AMBER),
    ("🔍", "Transaction Search",    '"Is my salary credited?" → searches by keyword → speaks result with date and amount.',          ORANGE),
    ("📘", "Cheque Book Request",   "Registers request, speaks delivery timeframe, then asks for next action.",                       BLUE),
    ("🔒", "Auto Biometric",        "Fingerprint auth triggers automatically — no button press. mic bar hides during scan.",          TEAL),
    ("📍", "Fixed Mic Button",      "Mic is pinned at the bottom of every screen. Never scrolls away. Always within reach.",          AMBER),
    ("⏱️",  "Session Timeout",       "3-min inactivity → spoken 30s warning → auto-return to auth. No data left exposed.",          ORANGE),
    ("🌐", "Hindi + English",       "Full bilingual support. Voice recognition and TTS both adapt to selected language.",            BLUE),
    ("🔠", "Font Scaling",          "5 steps from Smaller to Largest in Settings. Scales all text app-wide instantly.",              TEAL),
    ("🎨", "Colour-Blind Safe",     "Teal for success, orange for error, amber for accent. No red/green distinction anywhere.",      AMBER),
    ("💬", "Conversational AI",     "Every action ends with a follow-up question. No need to memorise commands.",                    ORANGE),
]
cols, rows = 4, 3
fw = Inches(3.18)
fh = Inches(1.56)
fx0 = Inches(0.4)
fy0 = Inches(1.85)
gap_x = Inches(0.1)
gap_y = Inches(0.12)

for idx, (icon, title, desc, clr) in enumerate(features):
    r, c = divmod(idx, cols)
    x = fx0 + (fw + gap_x) * c
    y = fy0 + (fh + gap_y) * r
    card_rect(sl, x, y, fw, fh, accent_color=clr)
    add_textbox(sl, icon, x + Inches(0.13), y + Inches(0.1), Inches(0.55), Inches(0.55), font_size=20)
    add_textbox(sl, title, x + Inches(0.72), y + Inches(0.1), fw - Inches(0.82), Inches(0.45),
                font_size=13, bold=True, color=WHITE)
    add_textbox(sl, desc,  x + Inches(0.13), y + Inches(0.65), fw - Inches(0.22), Inches(0.85),
                font_size=10.5, color=MUTED, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ACCESSIBILITY DESIGN
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "ACCESSIBILITY DESIGN")
slide_title(sl, "WCAG AAA — the gold standard")
accent_bar(sl)

# Left: donut/pie chart — WCAG compliance areas
cd2 = ChartData()
cd2.categories = ["Contrast", "Touch Targets", "Screen Reader", "Colour-Blind Safe", "Font Scale"]
cd2.add_series("Score", (20, 20, 20, 20, 20))

chart2 = sl.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT,
    Inches(0.3), Inches(1.88), Inches(4.5), Inches(4.5), cd2
).chart
chart2.has_legend  = False
chart2.has_title   = False
slice_colors = [BLUE, TEAL, AMBER, ORANGE, RGBColor(0xAB,0x47,0xBC)]
for i, pt in enumerate(chart2.plots[0].series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = slice_colors[i]

# Legend labels
legends = [
    ("Contrast  18.7:1", BLUE),
    ("Touch Targets  88dp", TEAL),
    ("Screen Reader  Full ARIA", AMBER),
    ("Colour-Blind Safe  Teal+Orange", ORANGE),
    ("Font Scale  5 steps", RGBColor(0xAB,0x47,0xBC)),
]
for i, (lbl, clr) in enumerate(legends):
    ly = Inches(2.35) + Inches(0.56) * i
    add_rect(sl, Inches(4.68), ly, Inches(0.22), Inches(0.22), fill=clr)
    add_textbox(sl, lbl, Inches(4.98), ly - Inches(0.02), Inches(2.3), Inches(0.34),
                font_size=12, color=WHITE)

# Right: comparison table
tx = Inches(7.3)
add_textbox(sl, "Our Implementation vs. Standard Apps",
            tx, Inches(1.82), Inches(5.7), Inches(0.45), font_size=14, bold=True, color=WHITE)

rows_data = [
    ("Contrast ratio",    "Standard: 4.5:1 (AA)",    "Saral: 18.7:1 (AAA ✓)", TEAL),
    ("Font size (min)",   "Standard: 12–14sp",        "Saral: 18sp base ✓",    TEAL),
    ("Tap target",        "Standard: 44dp",            "Saral: 88dp (2×) ✓",   TEAL),
    ("Colour dependency", "Standard: Red = error",     "Saral: Shape + icon ✓", TEAL),
    ("Screen reader",     "Standard: Partial",         "Saral: Full ARIA ✓",   TEAL),
    ("Mic button",        "Standard: Scrolls away",    "Saral: Fixed bottom ✓", TEAL),
    ("Biometric",         "Standard: Button tap",      "Saral: Auto-trigger ✓", TEAL),
]
for i, (req, std, saral, clr) in enumerate(rows_data):
    ry = Inches(2.35) + Inches(0.6) * i
    bg = NAVY_CARD if i % 2 == 0 else RGBColor(0x10, 0x24, 0x3C)
    add_rect(sl, tx, ry, Inches(5.72), Inches(0.56), fill=bg)
    add_textbox(sl, req,   tx + Inches(0.12), ry + Inches(0.1), Inches(1.6),  Inches(0.38), font_size=11, bold=True, color=WHITE)
    add_textbox(sl, std,   tx + Inches(1.75), ry + Inches(0.1), Inches(1.9),  Inches(0.38), font_size=11, color=ORANGE)
    add_textbox(sl, saral, tx + Inches(3.7),  ry + Inches(0.1), Inches(1.95), Inches(0.38), font_size=11, color=clr)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — COLOUR PALETTE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "DESIGN SYSTEM")
slide_title(sl, "Colour-blind safe palette — no red/green distinction")
accent_bar(sl)

palette = [
    ("#081320", NAVY,                  "Background",  "Navy Dark",          "Base canvas"),
    ("#4D9EFF", BLUE,                  "Primary",     "Accessible Blue",    "Buttons, links, focus rings"),
    ("#1556B0", RGBColor(0x15,0x56,0xB0), "Button Fill", "Deep Blue",       "7.1:1 contrast on white text"),
    ("#00BFA5", TEAL,                  "Success",     "Teal",               "Deuteranopia-safe ✓"),
    ("#FF6B35", ORANGE,                "Error",       "Deep Orange",        "Protanopia-safe ✓"),
    ("#FFC107", AMBER,                 "Accent",      "Amber",              "11.5:1 on dark bg (AAA)"),
    ("#FFFFFF", WHITE,                 "Text Primary","Pure White",          "On navy: 18.7:1 (AAA)"),
    ("#B2CCE8", MUTED,                 "Text Muted",  "Cool Blue-grey",     "Secondary info text"),
]
sw2 = Inches(1.52)
for i, (hex_val, clr, role, name, note) in enumerate(palette):
    x = Inches(0.4) + (sw2 + Inches(0.08)) * i
    # Swatch
    add_rect(sl, x, Inches(2.0), sw2, Inches(1.8), fill=clr,
             line=NAVY_LIGHT, line_w=Pt(1))
    if clr == WHITE:
        add_textbox(sl, hex_val, x, Inches(2.5), sw2, Inches(0.5),
                    font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    else:
        add_textbox(sl, hex_val, x, Inches(2.5), sw2, Inches(0.5),
                    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Labels
    add_textbox(sl, role,    x, Inches(3.9),  sw2, Inches(0.38), font_size=11, bold=True, color=WHITE,  align=PP_ALIGN.CENTER)
    add_textbox(sl, name,    x, Inches(4.26), sw2, Inches(0.38), font_size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(sl, note,    x, Inches(4.62), sw2, Inches(0.7),  font_size=9.5, color=MUTED, align=PP_ALIGN.CENTER, wrap=True, italic=True)

# Colour-blindness simulation note
add_rect(sl, Inches(0.4), Inches(5.55), Inches(12.4), Inches(1.72), fill=RGBColor(0x0A,0x20,0x38))
add_rect(sl, Inches(0.4), Inches(5.55), Inches(0.08), Inches(1.72), fill=AMBER)
add_textbox(sl, "Why Teal + Orange instead of Green + Red?",
            Inches(0.65), Inches(5.65), Inches(12.0), Inches(0.4), font_size=14, bold=True, color=AMBER)
add_textbox(sl,
    "Deuteranopia (green-blind, ~4.6% of males — Colour Blind Awareness): cannot distinguish green from red. "
    "Teal (#00BFA5) reads clearly as a distinct hue to deuteranopes.  |  "
    "Protanopia (red-blind, ~1% of males): cannot see red hues. "
    "Deep orange (#FF6B35) is distinguishable across all colour-blindness types.  |  "
    "Amber (#FFC107) achieves 11.5:1 contrast on navy (WebAIM calculator) — exceeding WCAG AAA (7:1).",
    Inches(0.65), Inches(6.08), Inches(12.0), Inches(1.1),
    font_size=12, color=WHITE, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECHNICAL ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "TECHNICAL ARCHITECTURE")
slide_title(sl, "How Saral is built")
accent_bar(sl)

# Android column
ax = Inches(0.4)
card_rect(sl, ax, Inches(1.82), Inches(6.05), Inches(5.45))
add_textbox(sl, "📱  Android (Kotlin + Compose)", ax + Inches(0.18), Inches(1.92),
            Inches(5.7), Inches(0.45), font_size=15, bold=True, color=BLUE)
add_rect(sl, ax + Inches(0.18), Inches(2.42), Inches(5.6), Inches(0.045), fill=NAVY_LIGHT)

android_lines = [
    "MVVM  +  Hilt Dependency Injection",
    "StateFlow / SharedFlow for UI state",
    "BiometricPrompt  →  auto-triggered via LaunchedEffect(step)",
    "SpeechRecognizer API  (voice input)",
    "TextToSpeech API  (voice output)",
    "Compose Box layout  →  fixed bottom mic bar",
    "CompositionLocalProvider(LocalDensity)  →  runtime font scale",
    "Synchronous callback pattern  →  testable side-effects",
]
for i, line in enumerate(android_lines):
    ly = Inches(2.57) + Inches(0.52) * i
    add_rect(sl, ax + Inches(0.22), ly + Inches(0.14), Inches(0.12), Inches(0.12), fill=BLUE)
    add_textbox(sl, line, ax + Inches(0.44), ly, Inches(5.5), Inches(0.46),
                font_size=11.5, color=WHITE)

# Web column
wx = Inches(6.7)
card_rect(sl, wx, Inches(1.82), Inches(6.2), Inches(5.45))
add_textbox(sl, "🌐  Web (HTML + CSS + Vanilla JS)", wx + Inches(0.18), Inches(1.92),
            Inches(5.9), Inches(0.45), font_size=15, bold=True, color=TEAL)
add_rect(sl, wx + Inches(0.18), Inches(2.42), Inches(5.8), Inches(0.045), fill=NAVY_LIGHT)

web_lines = [
    "Web Speech API  →  voice in + TTS out",
    "State machine  (currentTransferPhase, flags)",
    "awaitingSimpleFollowUp  →  post-action dialogue",
    "awaitingPostTransactionChoice  →  transaction follow-up",
    "CSS flex layout: scroll-content + mic-bar (fixed)",
    "showTransferBiometric()  →  setTimeout auto-execute (2.2s)",
    "Session timeout  →  clearTimeout / expireSession()",
    "Zoom  →  html { font-size }  + localStorage persistence",
]
for i, line in enumerate(web_lines):
    ly = Inches(2.57) + Inches(0.52) * i
    add_rect(sl, wx + Inches(0.22), ly + Inches(0.14), Inches(0.12), Inches(0.12), fill=TEAL)
    add_textbox(sl, line, wx + Inches(0.44), ly, Inches(5.7), Inches(0.46),
                font_size=11.5, color=WHITE)

# Test strip
add_rect(sl, Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.25), fill=RGBColor(0x05,0x12,0x22))
add_textbox(sl,
    "🧪  439 unit tests  •  0 failures  •  HomeViewModelTest  |  TransferViewModelTest  |  AuthViewModelTest  |  Domain / UseCases",
    Inches(0.55), Inches(7.13), Inches(12.2), Inches(0.24),
    font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONVERSATIONAL AI MODEL
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "VOICE AI")
slide_title(sl, "Multi-turn conversational dialogue")
accent_bar(sl)

add_textbox(sl,
    "Saral uses a flag-based intent state machine to maintain context across turns. "
    "Every action naturally leads into the next question — no command memorisation needed.",
    Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.65),
    font_size=14, color=MUTED, wrap=True)

# Dialogue examples — 3 columns
dialogues = [
    ("💳 Balance", BLUE, [
        ("User:", "Check my balance", WHITE),
        ("Saral:", "Your SBI account ending 7845 has ₹25,000 available.", TEAL),
        ("Saral:", "Say main menu to go home, or tell me what you'd like to do next.", MUTED),
        ("User:", "Transfer money", WHITE),
        ("Saral:", "[Navigates to transfer flow]", AMBER),
    ]),
    ("📋 Transactions", AMBER, [
        ("User:", "Last 5 transactions", WHITE),
        ("Saral:", "1. Salary ₹45,000 credited 25 May…  [reads all 5]", TEAL),
        ("Saral:", "Would you like to hear more, or say main menu?", MUTED),
        ("User:", "more", WHITE),
        ("Saral:", "Do you want last 1, last 5, or last 10?", TEAL),
    ]),
    ("💸 Transfer", TEAL, [
        ("User:", "Transfer 500 to Rahul", WHITE),
        ("Saral:", "You selected Rahul Sharma. Say the amount.", TEAL),
        ("User:", "five hundred", WHITE),
        ("Saral:", "Transfer ₹500 to Rahul. Say yes to confirm.", TEAL),
        ("User:", "yes", WHITE),
        ("Saral:", "[Fingerprint auto-scans → Transfer complete]", AMBER),
    ]),
]
for i, (title, clr, turns) in enumerate(dialogues):
    dx = Inches(0.4) + Inches(4.25) * i
    card_rect(sl, dx, Inches(2.55), Inches(4.07), Inches(4.7))
    add_rect(sl, dx, Inches(2.55), Inches(4.07), Inches(0.07), fill=clr)
    add_textbox(sl, title, dx + Inches(0.18), Inches(2.66), Inches(3.85), Inches(0.42),
                font_size=14, bold=True, color=WHITE)
    for j, (speaker, text, tcl) in enumerate(turns):
        ty = Inches(3.22) + Inches(0.71) * j
        add_textbox(sl, speaker, dx + Inches(0.18), ty, Inches(0.68), Inches(0.4),
                    font_size=11, bold=True, color=clr)
        add_textbox(sl, text, dx + Inches(0.88), ty, Inches(3.05), Inches(0.65),
                    font_size=11, color=tcl, wrap=True, italic=(speaker == "Saral:"))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TEST COVERAGE & RELEASE HISTORY
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "QUALITY & RELEASES")
slide_title(sl, "Built test-first — 439 tests, 0 failures")
accent_bar(sl)

# Test coverage bar chart
cd3 = ChartData()
cd3.categories = ["HomeViewModel", "TransferViewModel", "AuthViewModel", "Domain / UseCases"]
cd3.add_series("Tests", (55, 320, 18, 46))

chart3 = sl.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.35), Inches(1.85), Inches(5.6), Inches(4.0), cd3
).chart
chart3.has_legend = False
chart3.has_title  = False
bar_colors = [BLUE, TEAL, AMBER, ORANGE]
for i, pt in enumerate(chart3.plots[0].series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = bar_colors[i]
chart3.category_axis.tick_labels.font.size = Pt(11)
chart3.category_axis.tick_labels.font.color.rgb = MUTED
chart3.value_axis.tick_labels.font.size = Pt(11)
chart3.value_axis.tick_labels.font.color.rgb = MUTED

# Total stat
stat_block(sl, Inches(0.4), Inches(5.88), Inches(2.6), Inches(1.32), "439", "Total Tests", BLUE)
stat_block(sl, Inches(3.15), Inches(5.88), Inches(2.6), Inches(1.32), "0", "Failures", TEAL)

# Release timeline right
rx2 = Inches(6.2)
add_textbox(sl, "Release History", rx2, Inches(1.85), Inches(6.8), Inches(0.42),
            font_size=14, bold=True, color=WHITE)

releases = [
    ("v1.0", "Initial voice banking — balance, transfer, transactions",     MUTED),
    ("v3.0", "Navigation + SaralNavGraph wiring",                           MUTED),
    ("v5.0", "Biometric authentication added (button-triggered)",           MUTED),
    ("v6.0", "WCAG redesign · session timeout · beneficiary tap · logo",    BLUE),
    ("v7.0", "Balance + cheque book post-action follow-ups · 433 tests",    TEAL),
    ("v8.0", "Fixed mic bar · auto-trigger biometric · prompt doc · PPT",   AMBER),
]
for i, (ver, desc, clr) in enumerate(releases):
    ry2 = Inches(2.38) + Inches(0.82) * i
    is_current = (ver == "v8.0")
    bg = RGBColor(0x10, 0x24, 0x3C) if is_current else NAVY_CARD
    add_rect(sl, rx2, ry2, Inches(6.8), Inches(0.72), fill=bg,
             line=clr if is_current else NAVY_LIGHT, line_w=Pt(2 if is_current else 1))
    add_textbox(sl, ver,  rx2 + Inches(0.14), ry2 + Inches(0.12), Inches(0.65), Inches(0.46),
                font_size=13, bold=True, color=clr)
    add_textbox(sl, desc, rx2 + Inches(0.85), ry2 + Inches(0.12), Inches(5.8), Inches(0.52),
                font_size=11.5, color=WHITE if is_current else MUTED, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — IMPACT
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
section_label(sl, "IMPACT")
slide_title(sl, "Restoring financial independence at scale")
accent_bar(sl)

# Impact metrics pie chart (coverage of user groups)
cd4 = ChartData()
cd4.categories = ["Blind Users", "Low Vision", "Senior Citizens", "Colour-Blind", "Rural / Low Literacy"]
cd4.add_series("Impact (M)", (5, 75, 300, 300, 400))

chart4 = sl.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    Inches(0.3), Inches(1.9), Inches(5.2), Inches(4.5), cd4
).chart
chart4.has_legend = True
chart4.has_title  = False
pie_colors = [BLUE, TEAL, AMBER, ORANGE, RGBColor(0xAB,0x47,0xBC)]
for i, pt in enumerate(chart4.plots[0].series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = pie_colors[i]
legend = chart4.legend
legend.font.size  = Pt(10)
legend.font.color.rgb = MUTED

# Impact cards right
impact_items = [
    ("100%",  "of banking actions operable by voice alone",           BLUE),
    ("100%",  "of screens with fixed mic — always reachable",         TEAL),
    ("18.7:1","contrast ratio — exceeds WCAG AAA (7:1 minimum)",      AMBER),
    ("88dp",  "minimum tap target — double the Android specification", ORANGE),
    ("0",     "button taps required for fingerprint authentication",   TEAL),
    ("439",   "unit tests passing — zero failures",                   BLUE),
]
rx3 = Inches(5.75)
for i, (num, label, clr) in enumerate(impact_items):
    iy = Inches(1.87) + Inches(0.88) * i
    card_rect(sl, rx3, iy, Inches(7.15), Inches(0.76), accent_color=clr)
    add_textbox(sl, num,   rx3 + Inches(0.2),  iy + Inches(0.1), Inches(1.4), Inches(0.52),
                font_size=22, bold=True, color=clr, align=PP_ALIGN.RIGHT)
    add_textbox(sl, label, rx3 + Inches(1.7), iy + Inches(0.15), Inches(5.3), Inches(0.52),
                font_size=13, color=WHITE, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)

# Left accent panel
add_rect(sl, 0, 0, Inches(0.14), H, fill=BLUE)
add_rect(sl, Inches(0.14), 0, Inches(0.07), H, fill=TEAL)

# Big logo text
add_textbox(sl, "🏦🎙️", Inches(1.2), Inches(0.85), Inches(4), Inches(1.3), font_size=66, align=PP_ALIGN.CENTER)
add_textbox(sl, "SARAL", Inches(0.6), Inches(1.95), Inches(5.5), Inches(1.3),
            font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(sl, "सरल  —  Simple", Inches(0.6), Inches(3.1), Inches(5.5), Inches(0.55),
            font_size=20, color=AMBER, align=PP_ALIGN.CENTER, italic=True)
add_rect(sl, Inches(1.0), Inches(3.78), Inches(4.6), Inches(0.07), fill=BLUE)
add_textbox(sl, "Voice-First Accessible Banking", Inches(0.6), Inches(3.9), Inches(5.5), Inches(0.5),
            font_size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Right: quote + stats
add_textbox(sl, "\"Every person deserves the dignity of\nmanaging their own money —\nregardless of sight, age, or literacy.\"",
            Inches(6.1), Inches(1.3), Inches(6.9), Inches(1.85),
            font_size=21, bold=True, color=WHITE, wrap=True, italic=True)
add_rect(sl, Inches(6.1), Inches(3.22), Inches(6.9), Inches(0.07), fill=AMBER)

# 4 stats
stats_cl = [
    ("439",  "Tests Passing",        BLUE),
    ("0",    "Test Failures",        TEAL),
    ("AAA",  "WCAG Level",           AMBER),
    ("v8.0", "Current Release",      ORANGE),
]
sw3 = Inches(1.62)
for i, (num, lbl, clr) in enumerate(stats_cl):
    sx = Inches(6.1) + (sw3 + Inches(0.06)) * i
    stat_block(sl, sx, Inches(3.42), sw3, Inches(1.4), num, lbl, clr)

add_textbox(sl, "Android  •  Web  •  Hindi + English  •  ESFAI Hackathon 2026",
            Inches(6.1), Inches(5.0), Inches(6.9), Inches(0.42),
            font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

tag_data = [
    ("Android APK", BLUE,                    RGBColor(0x00,0x30,0x80)),
    ("Web App",     TEAL,                    RGBColor(0x00,0x60,0x55)),
    ("WCAG AAA",    AMBER,                   RGBColor(0x80,0x60,0x00)),
    ("Voice AI",    ORANGE,                  RGBColor(0x80,0x30,0x00)),
    ("Open Source", RGBColor(0xAB,0x47,0xBC),RGBColor(0x55,0x20,0x60)),
]
tx2 = Inches(6.1)
for j, (tag, clr, bg) in enumerate(tag_data):
    add_rect(sl, tx2, Inches(5.62), Inches(1.3), Inches(0.38),
             fill=bg, line=clr, line_w=Pt(1.5))
    add_textbox(sl, tag, tx2 + Inches(0.06), Inches(5.64), Inches(1.18), Inches(0.34),
                font_size=10, bold=True, color=clr, align=PP_ALIGN.CENTER)
    tx2 += Inches(1.38)

# ESFAI branding bottom
add_rect(sl, 0, Inches(7.1), W, Inches(0.4), fill=RGBColor(0x04,0x0C,0x18))
add_textbox(sl, "ESFAI Hackathon 2026  •  Team Saral  •  esfaihackathon@gmail.com",
            Inches(0.3), Inches(7.12), W - Inches(0.6), Inches(0.32),
            font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — REFERENCES
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_bg(sl)
add_rect(sl, 0, 0, Inches(0.14), H, fill=AMBER)
section_label(sl, "REFERENCES & DATA SOURCES")
slide_title(sl, "All data cited — sources and accuracy notes")
accent_bar(sl)

refs = [
    # (category_label, category_color, [(citation_text, note)])
    ("Vision Impairment Data", BLUE, [
        ("[1] IAPB Vision Atlas 2023 — iapb.org/learn/iapb-vision-atlas",
         "India: 8.95M blind, 62.7M moderate-severe vision impairment = ~72M total. Used on Slides 2, 3."),
        ("[2] WHO World Report on Vision, 2019 — WHO/NMH/PBD/19.1",
         "Global overview of vision impairment burden; India is the country with the highest number of blind people globally."),
    ]),
    ("Ageing & Senior Citizens", TEAL, [
        ("[3] UN World Population Ageing 2023 — un.org/development/desa/pd",
         "India 60+ population: ~145M current (2023); projected 347M by 2050. Slide 3 uses current figure with projection noted."),
        ("[4] India Ageing Report 2023 — MoSPI / UNFPA India",
         "Published by Ministry of Statistics & Programme Implementation; confirms UN projections for India's elderly population trajectory."),
    ]),
    ("Colour-Blindness", ORANGE, [
        ("[5] Colour Blind Awareness — colourblindawareness.org",
         "~300M people globally are colour-blind; ~8% of males and 0.5% of females affected (all types combined). India ≈ 60M."),
        ("[6] Deeb, S.S. (2004). The molecular basis of variation in human color vision. Clinical Genetics, 67(5), 369–377.",
         "Peer-reviewed basis for deuteranopia (~4.63% males) and protanopia (~1.01% males) prevalence figures."),
    ]),
    ("Literacy & Rural Access", AMBER, [
        ("[7] Census of India 2011 — censusindia.gov.in",
         "Literacy rate 74.04%. Applied to 2023 population (~1.43B) gives ~372M with limited literacy. Used on Slides 2, 3."),
        ("[8] National Literacy Mission, Govt. of India — nlm.nic.in",
         "Functional literacy definition and adult literacy programme data."),
    ]),
    ("Accessibility Standards", RGBColor(0xAB,0x47,0xBC), [
        ("[9] W3C WCAG 2.1 Guidelines — w3.org/TR/WCAG21",
         "SC 1.4.3: AA contrast 4.5:1 normal text. SC 1.4.6: AAA contrast 7:1. SC 1.4.4: text resizable to 200%. Slides 7, 8."),
        ("[10] WebAIM Contrast Checker — webaim.org/resources/contrastchecker",
         "Verified contrast ratios: White #FFF on Navy #081320 = 18.7:1 (AAA). Amber #FFC107 on Navy = 11.5:1 (AAA)."),
        ("[11] Material Design Accessibility — m3.material.io/foundations/accessible-design",
         "Recommended minimum touch target: 48×48dp. Saral uses 88dp (1.83× the recommendation). Slide 7."),
    ]),
]

row_y = Inches(1.82)
for cat_label, cat_color, items in refs:
    # Category header
    add_rect(sl, Inches(0.4), row_y, Inches(12.5), Inches(0.36), fill=RGBColor(0x0A,0x1E,0x34))
    add_rect(sl, Inches(0.4), row_y, Inches(0.08), Inches(0.36), fill=cat_color)
    add_textbox(sl, cat_label, Inches(0.6), row_y + Inches(0.04), Inches(12.0), Inches(0.3),
                font_size=12, bold=True, color=cat_color)
    row_y += Inches(0.38)
    for cite, note in items:
        add_rect(sl, Inches(0.4), row_y, Inches(12.5), Inches(0.54), fill=NAVY_CARD)
        add_textbox(sl, cite, Inches(0.55), row_y + Inches(0.03), Inches(12.1), Inches(0.25),
                    font_size=10.5, bold=True, color=WHITE)
        add_textbox(sl, note, Inches(0.55), row_y + Inches(0.27), Inches(12.1), Inches(0.24),
                    font_size=9.5, color=MUTED, italic=True, wrap=True)
        row_y += Inches(0.56)
    row_y += Inches(0.04)  # gap between categories

# Bottom disclaimer
add_rect(sl, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.28), fill=RGBColor(0x04,0x0C,0x18))
add_textbox(sl,
    "All population figures are the best publicly available estimates at time of presentation. "
    "India-specific and global figures are distinguished in slide notes. Contrast ratios computed per WCAG 2.1 relative luminance formula.",
    Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.24),
    font_size=9, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = "/Users/aman/Desktop/esfaihackathon/visual-assistant-ai/presentation/Saral-Accessible-Banking.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
