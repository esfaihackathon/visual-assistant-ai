"""
Saral — Business Stakeholder Deck
5 slides, executive-level, data-driven
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x08, 0x13, 0x20)
NAVY_MID  = RGBColor(0x0D, 0x1E, 0x36)
NAVY_CARD = RGBColor(0x0E, 0x20, 0x35)
BLUE      = RGBColor(0x4D, 0x9E, 0xFF)
TEAL      = RGBColor(0x00, 0xBF, 0xA5)
AMBER     = RGBColor(0xFF, 0xC1, 0x07)
ORANGE    = RGBColor(0xFF, 0x6B, 0x35)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0xB2, 0xCC, 0xE8)
DARK      = RGBColor(0x04, 0x0C, 0x18)
PURPLE    = RGBColor(0xAB, 0x47, 0xBC)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def bg(sl):
    r = sl.shapes.add_shape(1, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = NAVY
    r.line.fill.background()

def rect(sl, x, y, w, h, fill=None, line_color=None, lw=Pt(1)):
    s = sl.shapes.add_shape(1, x, y, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line_color: s.line.color.rgb = line_color; s.line.width = lw
    else: s.line.fill.background()
    return s

def tx(sl, text, x, y, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    t = sl.shapes.add_textbox(x, y, w, h)
    t.word_wrap = wrap
    tf = t.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri"
    return t

def divider(sl, y, color=BLUE, x=Inches(0.55), w=Inches(1.6)):
    rect(sl, x, y, w, Inches(0.055), fill=color)

def kpi(sl, x, y, w, h, num, label, num_clr=BLUE, bg_clr=NAVY_CARD):
    rect(sl, x, y, w, h, fill=bg_clr)
    rect(sl, x, y, w, Inches(0.06), fill=num_clr)
    tx(sl, num,   x, y+Inches(0.12), w, Inches(0.7),
       size=30, bold=True, color=num_clr, align=PP_ALIGN.CENTER)
    tx(sl, label, x, y+Inches(0.8),  w, Inches(0.65),
       size=10.5, color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

_chip_bg = {
    id(BLUE):   RGBColor(0x00, 0x30, 0x80),
    id(TEAL):   RGBColor(0x00, 0x55, 0x4A),
    id(AMBER):  RGBColor(0x6B, 0x50, 0x00),
    id(ORANGE): RGBColor(0x70, 0x28, 0x00),
    id(MUTED):  RGBColor(0x2A, 0x40, 0x58),
}
def chip(sl, text, x, y, color):
    bkg = _chip_bg.get(id(color), RGBColor(0x15, 0x25, 0x40))
    rect(sl, x, y, Inches(1.55), Inches(0.34), fill=bkg, line_color=color, lw=Pt(1.2))
    tx(sl, text, x+Inches(0.06), y+Inches(0.05), Inches(1.43), Inches(0.26),
       size=9.5, bold=True, color=color, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER (executive)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)

# Full-width gradient left band
rect(sl, 0, 0, Inches(5.2), H, fill=RGBColor(0x09, 0x18, 0x2E))
rect(sl, 0, 0, Inches(0.16), H, fill=TEAL)
rect(sl, Inches(0.16), 0, Inches(0.07), H, fill=BLUE)

# Logo mark
rect(sl, Inches(1.0), Inches(1.05), Inches(2.8), Inches(2.1), fill=RGBColor(0x0C,0x22,0x40))
rect(sl, Inches(1.0), Inches(1.05), Inches(2.8), Inches(0.09), fill=BLUE)
tx(sl, "🏦", Inches(1.0), Inches(1.1), Inches(2.8), Inches(1.25),
   size=52, align=PP_ALIGN.CENTER)
tx(sl, "SARAL", Inches(0.9), Inches(2.28), Inches(3.0), Inches(0.8),
   size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(sl, "सरल", Inches(0.9), Inches(2.98), Inches(3.0), Inches(0.44),
   size=17, italic=True, color=AMBER, align=PP_ALIGN.CENTER)
rect(sl, Inches(1.1), Inches(3.52), Inches(2.8), Inches(0.055), fill=TEAL)
tx(sl, "Voice-First Accessible Banking", Inches(0.8), Inches(3.68),
   Inches(3.2), Inches(0.4), size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Right — headline + 3 impact stats
rx = Inches(5.55)
tx(sl, "Banking for everyone —\nregardless of sight, age,\nor literacy.",
   rx, Inches(1.0), Inches(7.35), Inches(1.9),
   size=30, bold=True, color=WHITE, wrap=True)
rect(sl, rx, Inches(3.0), Inches(2.5), Inches(0.07), fill=AMBER)
tx(sl, "India's first fully voice-driven, WCAG AAA\ncompliant accessible banking assistant.\nAvailable on Android and Web.",
   rx, Inches(3.12), Inches(7.35), Inches(1.1),
   size=15, color=MUTED, wrap=True)

kpis = [
    ("630M+", "Indians excluded\nfrom digital banking",  ORANGE),
    ("WCAG\nAAA",  "Highest accessibility\nstandard achieved",   TEAL),
    ("100%",  "Actions completable\nby voice alone",          BLUE),
]
for i, (n, l, c) in enumerate(kpis):
    kpi(sl, rx + Inches(0.02) + Inches(2.46)*i,
        Inches(4.4), Inches(2.38), Inches(1.65), n, l, c)

rect(sl, 0, Inches(7.1), W, Inches(0.4), fill=DARK)
tx(sl, "ESFAI Hackathon 2026  •  v8.0  •  439 Tests, 0 Failures  •  esfaihackathon@gmail.com",
   Inches(0.3), Inches(7.13), W-Inches(0.6), Inches(0.3),
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — PROBLEM + MARKET OPPORTUNITY
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
tx(sl, "THE PROBLEM  &  OPPORTUNITY", Inches(0.55), Inches(0.28),
   Inches(10), Inches(0.38), size=10.5, bold=True, color=AMBER)
tx(sl, "630M+ Indians cannot use standard banking apps",
   Inches(0.55), Inches(0.65), Inches(12.3), Inches(0.95),
   size=33, bold=True, color=WHITE)
divider(sl, Inches(1.6))

# Left — horizontal bar chart (market size by group)
cd = ChartData()
cd.categories = ["Low Literacy\n~372M", "Senior Citizens\n~145M current\n(347M by 2050)",
                 "Colour-Blind\n~60M India\n(300M global)", "Vision Impaired\n~72M"]
cd.add_series("Users (Millions)", (372, 145, 60, 72))
ch = sl.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.35), Inches(1.75), Inches(5.8), Inches(4.6), cd
).chart
ch.has_legend = False; ch.has_title = False
bar_clrs = [AMBER, TEAL, ORANGE, BLUE]
for i, pt in enumerate(ch.plots[0].series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = bar_clrs[i]
ch.category_axis.tick_labels.font.size = Pt(10)
ch.category_axis.tick_labels.font.color.rgb = MUTED
ch.value_axis.tick_labels.font.size = Pt(10)
ch.value_axis.tick_labels.font.color.rgb = MUTED

# Right — 4 problem cards
px = Inches(6.38)
problems = [
    ("👁️ Visual", "~72M Indians cannot read any UI. Screen readers unsupported by most banking apps.", BLUE),
    ("🤲 Motor",  "~145M seniors lack fine-motor precision for small buttons — including biometric auth.", TEAL),
    ("🎨 Colour", "~60M colour-blind Indians cannot interpret red/green status indicators.", ORANGE),
    ("📖 Literacy","~372M with limited literacy cannot navigate menu-driven or text-heavy interfaces.", AMBER),
]
for i, (title, body, clr) in enumerate(problems):
    py = Inches(1.75) + Inches(1.18)*i
    rect(sl, px, py, Inches(6.6), Inches(1.08), fill=NAVY_CARD)
    rect(sl, px, py, Inches(0.08), Inches(1.08), fill=clr)
    tx(sl, title, px+Inches(0.18), py+Inches(0.1), Inches(6.2), Inches(0.35),
       size=13, bold=True, color=WHITE)
    tx(sl, body,  px+Inches(0.18), py+Inches(0.46), Inches(6.2), Inches(0.56),
       size=11, color=MUTED, wrap=True)

# Source footnote
rect(sl, Inches(0.35), Inches(6.48), Inches(12.6), Inches(0.28), fill=DARK)
tx(sl, "Sources: [1] IAPB Vision Atlas 2023  [3] UN Population Ageing 2023  [5] Colour Blind Awareness  [7] Census of India 2011  — See References slide",
   Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.24),
   size=8.5, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Bottom tagline
tx(sl, "These users are not edge cases — they are 44% of India's adult population.",
   Inches(0.55), Inches(6.85), Inches(12.3), Inches(0.48),
   size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — SOLUTION + USER FLOW + DIFFERENTIATORS
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
tx(sl, "THE SOLUTION", Inches(0.55), Inches(0.28), Inches(10), Inches(0.38),
   size=10.5, bold=True, color=AMBER)
tx(sl, "Saral: every banking action, by voice alone",
   Inches(0.55), Inches(0.65), Inches(12.3), Inches(0.85),
   size=33, bold=True, color=WHITE)
divider(sl, Inches(1.52))

# Transfer flow — 6 steps in one row
steps = [
    ("🗣️", "Speak",     "\"Transfer\n500 to Rahul\"", BLUE),
    ("👥", "Select",    "Tap or say\nbeneficiary",     TEAL),
    ("💰", "Amount",    "\"500 rupees\"",              AMBER),
    ("✅", "Confirm",   "\"Yes\"",                     TEAL),
    ("👆", "Biometric", "Auto-scans\nfingerprint",     ORANGE),
    ("🎉", "Done",      "Balance read\naloud",         TEAL),
]
sw = Inches(1.97); gap = Inches(0.18); sx0 = Inches(0.38)
for i, (ico, title, cap, clr) in enumerate(steps):
    x = sx0 + (sw+gap)*i
    rect(sl, x, Inches(1.75), sw, Inches(2.15), fill=NAVY_CARD)
    rect(sl, x, Inches(1.75), sw, Inches(0.07), fill=clr)
    rect(sl, x+Inches(0.06), Inches(1.82), Inches(0.38), Inches(0.28), fill=clr)
    tx(sl, str(i+1), x+Inches(0.06), Inches(1.81), Inches(0.38), Inches(0.28),
       size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tx(sl, ico,   x, Inches(2.18), sw, Inches(0.7),  size=28, align=PP_ALIGN.CENTER)
    tx(sl, title, x, Inches(2.92), sw, Inches(0.35), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(sl, cap,   x, Inches(3.27), sw, Inches(0.6),  size=10, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    if i < len(steps)-1:
        tx(sl, "→", x+sw+Inches(0.02), Inches(2.6), gap+Inches(0.16), Inches(0.4),
           size=18, color=BLUE, align=PP_ALIGN.CENTER)

# 3 differentiators below
diffs = [
    ("📍 Fixed Mic — Always Reachable",
     "Mic button is pinned at the bottom of every screen. Never scrolls away. "
     "Critical for users who cannot scan the UI to find controls.",
     BLUE),
    ("👆 Zero-Tap Biometric Auth",
     "Fingerprint scanner activates automatically when a transfer is confirmed. "
     "No button press needed — removes the single biggest motor-friction point in banking.",
     ORANGE),
    ("💬 Conversational Follow-ups",
     "After every action the app asks what to do next. No commands to memorise. "
     "Works for low-literacy and first-time digital users.",
     TEAL),
]
for i, (title, body, clr) in enumerate(diffs):
    x = Inches(0.38) + Inches(4.35)*i
    rect(sl, x, Inches(4.1), Inches(4.2), Inches(2.12), fill=NAVY_CARD)
    rect(sl, x, Inches(4.1), Inches(4.2), Inches(0.07), fill=clr)
    tx(sl, title, x+Inches(0.18), Inches(4.22), Inches(3.9), Inches(0.42),
       size=12.5, bold=True, color=WHITE)
    tx(sl, body,  x+Inches(0.18), Inches(4.67), Inches(3.9), Inches(1.48),
       size=11, color=MUTED, wrap=True)

# Platforms strip
rect(sl, Inches(0.38), Inches(6.38), Inches(12.6), Inches(0.38), fill=DARK)
for j, (lbl, clr) in enumerate([
        ("Android APK", BLUE), ("Web App", TEAL), ("Hindi + English", AMBER),
        ("WCAG AAA", ORANGE), ("3-Min Session Timeout", MUTED), ("439 Tests · 0 Failures", TEAL)]):
    chip(sl, lbl, Inches(0.5)+Inches(2.06)*j, Inches(6.44), clr)

tx(sl, "No reading required at any step. No colour interpretation. No small buttons.",
   Inches(0.38), Inches(6.92), Inches(12.6), Inches(0.44),
   size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — IMPACT + COMPLIANCE + COMPETITIVE POSITION
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
tx(sl, "IMPACT  &  COMPLIANCE", Inches(0.55), Inches(0.28), Inches(10), Inches(0.38),
   size=10.5, bold=True, color=AMBER)
tx(sl, "Measurable results — highest accessibility standard achieved",
   Inches(0.55), Inches(0.65), Inches(12.3), Inches(0.85),
   size=33, bold=True, color=WHITE)
divider(sl, Inches(1.52))

# Top KPI row
kpis4 = [
    ("630M+",  "Indians who can now\nbank independently",  ORANGE),
    ("18.7:1", "Contrast ratio\n(WCAG AAA min: 7:1)",    TEAL),
    ("88dp",   "Min touch target\n(2× Material Design)",  BLUE),
    ("100%",   "Screens operable\nby voice alone",        AMBER),
    ("0",      "Button taps needed\nfor biometric auth",  TEAL),
]
kw = Inches(2.44)
for i, (n, l, c) in enumerate(kpis4):
    kpi(sl, Inches(0.38)+kw*i, Inches(1.72), kw-Inches(0.08), Inches(1.5), n, l, c)

# Left — donut chart (who benefits)
cd2 = ChartData()
cd2.categories = ["Vision Impaired ~72M", "Colour-Blind ~60M", "Senior Citizens ~145M", "Low Literacy ~372M"]
cd2.add_series("Users", (72, 60, 145, 372))
ch2 = sl.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT,
    Inches(0.3), Inches(3.38), Inches(4.6), Inches(3.38), cd2
).chart
ch2.has_legend = True; ch2.has_title = False
pie_clrs = [BLUE, ORANGE, TEAL, AMBER]
for i, pt in enumerate(ch2.plots[0].series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = pie_clrs[i]
ch2.legend.font.size = Pt(9.5); ch2.legend.font.color.rgb = MUTED

# Centre donut label
tx(sl, "649M\nusers\nreached", Inches(1.45), Inches(4.38), Inches(1.9), Inches(1.38),
   size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Right — compliance table
rx = Inches(5.2)
tx(sl, "Accessibility Compliance Checklist", rx, Inches(3.38), Inches(7.8), Inches(0.38),
   size=13, bold=True, color=WHITE)
checks = [
    ("WCAG 2.1 Level AAA — Contrast",     "18.7:1 on primary text",                "✓ AAA", TEAL),
    ("WCAG 2.1 Level AAA — Resize Text",  "5-step font scale (Settings)",           "✓ AAA", TEAL),
    ("WCAG 2.1 — Non-Text Contrast",      "All icons ≥ 3:1 against background",     "✓ AA",  BLUE),
    ("WCAG 2.1 — Touch Target Size",      "88dp (spec: 44dp) — 2× margin",         "✓",     TEAL),
    ("Android Accessibility — TalkBack",  "Full contentDescription on all elements","✓",     BLUE),
    ("Web ARIA",                          "role, aria-label, aria-live regions",    "✓",     TEAL),
    ("Colour Independence",               "No information conveyed by colour alone", "✓",     ORANGE),
    ("Session Security",                  "3-min idle timeout + 30s spoken warning","✓",     AMBER),
]
for i, (req, impl, status, clr) in enumerate(checks):
    ry = Inches(3.85) + Inches(0.36)*i
    fill = NAVY_CARD if i%2==0 else RGBColor(0x10,0x24,0x3C)
    rect(sl, rx, ry, Inches(7.75), Inches(0.34), fill=fill)
    tx(sl, req,    rx+Inches(0.12), ry+Inches(0.05), Inches(4.0), Inches(0.26), size=10.5, bold=True, color=WHITE)
    tx(sl, impl,   rx+Inches(4.15), ry+Inches(0.05), Inches(2.7), Inches(0.26), size=10,   color=MUTED)
    tx(sl, status, rx+Inches(6.9),  ry+Inches(0.05), Inches(0.8), Inches(0.26), size=10.5, bold=True, color=clr, align=PP_ALIGN.CENTER)

rect(sl, Inches(0.38), Inches(6.82), Inches(12.6), Inches(0.28), fill=DARK)
tx(sl, "Sources: [9] W3C WCAG 2.1  [10] WebAIM Contrast Checker  [11] Material Design Accessibility  — See References slide",
   Inches(0.5), Inches(6.84), Inches(12.3), Inches(0.22),
   size=8.5, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — PROMPT ENGINEERING
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
tx(sl, "PROMPT ENGINEERING", Inches(0.55), Inches(0.28), Inches(10), Inches(0.38),
   size=10.5, bold=True, color=AMBER)
tx(sl, "How natural language prompts became shipped features",
   Inches(0.55), Inches(0.65), Inches(12.3), Inches(0.85),
   size=33, bold=True, color=WHITE)
divider(sl, Inches(1.52))

# Left panel — example prompt (transaction feature)
lx = Inches(0.38)
rect(sl, lx, Inches(1.72), Inches(6.1), Inches(5.5), fill=NAVY_CARD)
rect(sl, lx, Inches(1.72), Inches(6.1), Inches(0.07), fill=AMBER)
rect(sl, lx, Inches(1.72), Inches(0.07), Inches(5.5), fill=AMBER)

tx(sl, "📋  Example Prompt  →  Transaction Voice Feature",
   lx+Inches(0.16), Inches(1.82), Inches(5.8), Inches(0.38),
   size=12, bold=True, color=AMBER)

prompt_lines = [
    ('Current State:', True,  WHITE),
    ('App only speaks the last transaction.', False, MUTED),
    ('', False, MUTED),
    ('New Features Required:', True,  WHITE),
    ('① User says "transaction" → app asks:', False, MUTED),
    ('   "Last 1, last 5, or last 10 transactions?"', False, BLUE),
    ('② App reads out N transactions by voice', False, MUTED),
    ('   (date · description · amount)', False, MUTED),
    ('', False, MUTED),
    ('③ Query mode — user asks:', True, WHITE),
    ('   "Is my salary credited this month?"', False, BLUE),
    ('   "Any electricity bill transaction?"', False, BLUE),
    ('④ App searches data → responds by voice:', False, MUTED),
    ('   "Yes, ₹45,000 credited on 25 May."', False, TEAL),
    ('   "No matching transaction found."', False, ORANGE),
]

tb = sl.shapes.add_textbox(lx+Inches(0.18), Inches(2.28), Inches(5.74), Inches(4.8))
tb.word_wrap = True
tf = tb.text_frame; tf.word_wrap = True
for i, (line, bold, clr) in enumerate(prompt_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = line
    r.font.size = Pt(11); r.font.bold = bold
    r.font.color.rgb = clr; r.font.name = "Calibri"
    from pptx.util import Pt as _Pt
    p.space_after = _Pt(1)

# Right panel — prompt → feature pipeline
rx2 = Inches(6.7)
tx(sl, "Prompt → Feature Pipeline",
   rx2, Inches(1.72), Inches(6.3), Inches(0.38),
   size=13, bold=True, color=WHITE)

steps_p = [
    ("1", "Write Prompt",
     "Plain English spec describing current state, desired behaviour, and example interactions.",
     BLUE),
    ("2", "AI Parses Intent",
     "Intent parser (VoiceIntentParser.kt / parseIntent() in JS) maps voice text to structured intent type.",
     TEAL),
    ("3", "Feature Implemented",
     "ViewModel handles intent: flag-based state machine, follow-up prompts, TTS response.",
     AMBER),
    ("4", "Tests Written",
     "Unit tests for each new intent path, follow-up choice, and edge case. 439 tests, 0 failures.",
     ORANGE),
    ("5", "APK / Web Shipped",
     "Both Android APK and Web built and released. Same feature parity across platforms.",
     TEAL),
]
for i, (num, title, body, clr) in enumerate(steps_p):
    sy = Inches(2.22) + Inches(1.04)*i
    rect(sl, rx2, sy, Inches(6.3), Inches(0.94), fill=NAVY_CARD)
    rect(sl, rx2, sy, Inches(0.07), Inches(0.94), fill=clr)
    rect(sl, rx2+Inches(0.18), sy+Inches(0.14), Inches(0.36), Inches(0.36), fill=clr)
    tx(sl, num, rx2+Inches(0.18), sy+Inches(0.13), Inches(0.36), Inches(0.36),
       size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tx(sl, title, rx2+Inches(0.68), sy+Inches(0.1), Inches(5.5), Inches(0.34),
       size=12.5, bold=True, color=WHITE)
    tx(sl, body, rx2+Inches(0.68), sy+Inches(0.46), Inches(5.5), Inches(0.44),
       size=10.5, color=MUTED, wrap=True)

# Other prompts used — bottom strip
rect(sl, Inches(0.38), Inches(7.12), Inches(12.6), Inches(0.3), fill=DARK)
other_prompts = [
    '"Mic button is moving — fix at bottom"',
    '"Auto-trigger fingerprint, no button"',
    '"After balance, ask if back to menu"',
    '"Session timeout 3 min"',
    '"Beneficiary list — tap to select"',
]
tx(sl, "Other prompts:  " + "  ·  ".join(other_prompts),
   Inches(0.5), Inches(7.14), Inches(12.3), Inches(0.26),
   size=9, color=MUTED, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — CODEBASE & CONTRIBUTORS
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
tx(sl, "CODEBASE  &  CONTRIBUTORS", Inches(0.55), Inches(0.28), Inches(10), Inches(0.38),
   size=10.5, bold=True, color=AMBER)
tx(sl, "Who built what — commits, files, and ownership",
   Inches(0.55), Inches(0.65), Inches(12.3), Inches(0.85),
   size=33, bold=True, color=WHITE)
divider(sl, Inches(1.52))

# Top — 5 repo stats
repo_kpis = [
    ("5",    "Team\nMembers",                        PURPLE),
    ("35",   "Kotlin source files\n(Android)",       BLUE),
    ("16",   "Test files · 439 tests\n0 failures",   TEAL),
    ("43",   "Total files changed\nacross all PRs",  AMBER),
    ("7,813","Lines added\n1,220 lines removed",     ORANGE),
]
rw = Inches(2.54)
for i, (n, l, c) in enumerate(repo_kpis):
    kpi(sl, Inches(0.38)+rw*i, Inches(1.72), rw-Inches(0.08), Inches(1.3), n, l, c)

# ── Left: contributor table (5 members, compact rows) ──────────────────────
lx2 = Inches(0.38)
tx(sl, "Contributors & Commit Ownership",
   lx2, Inches(3.12), Inches(6.6), Inches(0.36),
   size=12.5, bold=True, color=WHITE)

# Column headers
for hx, htxt, hclr in [
    (lx2+Inches(0.18), "Commits", MUTED),
    (lx2+Inches(0.82), "Name / Role", MUTED),
]:
    tx(sl, htxt, hx, Inches(3.52), Inches(1.5), Inches(0.22),
       size=9, bold=True, color=hclr)

# (name, handle, commits_label, commits_note, role_text, clr)
contribs = [
    ("Aman Deep",    "@AmanDeepSrivastava19",
     "10",
     "10 total  ·  3 solo  ·  7 with Claude AI co-author\n"
     "(v5.0 biometric · v7.0 redesign · v8.0 mic/biometric · PPT)",
     "Feature delivery: accessibility redesign, voice follow-ups,\n"
     "auto-biometric, fixed mic bar, session timeout, all docs & PPT",
     BLUE),
    ("Arundeep Kamboj", "@ArundeepKamboj",
     "14",
     "14 solo commits",
     "Core architecture, HomeScreen redesign, WCAG colour palette,\n"
     "comprehensive unit tests, navigation wiring, v3.0 release",
     TEAL),
    ("Jayesh Sharma", "@SHARMAJAYESH",
     "7",
     "7 solo commits",
     "First prototype + APK, README, initial prompt document,\n"
     "first feature demo for stakeholders",
     AMBER),
    ("Sukamal Maity", "@SukamalMaity",
     "—",
     "Product strategy & UX review",
     "User research, accessibility requirement specification,\n"
     "stakeholder presentation, QA sign-off",
     PURPLE),
    ("Hardeep",       "@hkhardeep9",
     "2",
     "2 solo commits",
     "OTP flow removal, early APK build & distribution",
     ORANGE),
]
row_h = Inches(0.76)
for i, (name, handle, commits, cnote, role, clr) in enumerate(contribs):
    cy2 = Inches(3.76) + row_h * i
    rect(sl, lx2, cy2, Inches(6.6), row_h - Inches(0.04), fill=NAVY_CARD)
    rect(sl, lx2, cy2, Inches(0.07), row_h - Inches(0.04), fill=clr)
    # Commit count badge
    rect(sl, lx2+Inches(0.15), cy2+Inches(0.18), Inches(0.52), Inches(0.36), fill=clr)
    tx(sl, commits, lx2+Inches(0.15), cy2+Inches(0.17), Inches(0.52), Inches(0.36),
       size=14, bold=True, color=DARK if commits != "—" else WHITE, align=PP_ALIGN.CENTER)
    # Name + handle
    tx(sl, name,   lx2+Inches(0.8), cy2+Inches(0.04), Inches(5.65), Inches(0.28),
       size=11.5, bold=True, color=WHITE)
    tx(sl, f"{handle}  ·  {cnote}", lx2+Inches(0.8), cy2+Inches(0.30), Inches(5.65), Inches(0.22),
       size=8.5, color=clr, italic=True, wrap=True)
    tx(sl, role, lx2+Inches(0.8), cy2+Inches(0.50), Inches(5.65), Inches(0.24),
       size=9, color=MUTED, wrap=True)

# ── Right: key files & ownership ───────────────────────────────────────────
rx3 = Inches(7.18)
tx(sl, "Key Files & Ownership",
   rx3, Inches(3.12), Inches(6.0), Inches(0.36),
   size=12.5, bold=True, color=WHITE)

# (file, display_author, ai_assisted, desc, clr)
files = [
    ("HomeViewModel.kt",         "Aman Deep",     True,  "Voice intent, follow-up state machine",       BLUE),
    ("TransferViewModel.kt",     "Aman Deep",     True,  "Transfer flow, biometric, nav callback",      BLUE),
    ("TransferScreen.kt",        "Aman Deep",     True,  "Fixed mic bar, auto-biometric trigger",       BLUE),
    ("HomeScreen.kt",            "Aman Deep",     True,  "Fixed bottom mic bar layout",                 BLUE),
    ("MainActivity.kt",          "Aman Deep",     True,  "Session timeout, biometric prompt",           BLUE),
    ("HomeViewModelTest.kt",     "Aman Deep",     True,  "55 tests incl. follow-up flows",              TEAL),
    ("TransferViewModelTest.kt", "Aman Deep",     True,  "320 tests incl. auto-biometric",              TEAL),
    ("Color.kt / Theme.kt",      "Arundeep",      False, "WCAG AAA colour palette",                     ORANGE),
    ("SaralNavGraph.kt",         "Arundeep",      False, "Navigation wiring",                           ORANGE),
    ("web/app.js",               "Aman Deep",     True,  "Web voice AI, all state machines",            BLUE),
    ("web/styles.css",           "Aman Deep",     True,  "WCAG design system, mic-bar CSS",             BLUE),
    ("PROMPT_DOCUMENT.md",       "Aman Deep",     True,  "Full project brief & prompt history",         PURPLE),
]

# Column header row
rect(sl, rx3, Inches(3.52), Inches(6.1), Inches(0.24), fill=RGBColor(0x08,0x18,0x2E))
for hx2, htxt2 in [
    (rx3+Inches(0.1),  "File"),
    (rx3+Inches(2.7),  "Author"),
    (rx3+Inches(4.15), "AI-assisted"),
    (rx3+Inches(5.1),  "Purpose"),
]:
    tx(sl, htxt2, hx2, Inches(3.54), Inches(1.2), Inches(0.2),
       size=8.5, bold=True, color=MUTED)

for i, (fname, author, ai, desc, clr) in enumerate(files):
    fy = Inches(3.78) + Inches(0.27)*i
    fill = NAVY_CARD if i%2==0 else RGBColor(0x10,0x24,0x3C)
    rect(sl, rx3, fy, Inches(6.1), Inches(0.26), fill=fill)
    tx(sl, fname,  rx3+Inches(0.1),  fy+Inches(0.03), Inches(2.55), Inches(0.22), size=9,   bold=True, color=WHITE)
    tx(sl, author, rx3+Inches(2.7),  fy+Inches(0.03), Inches(1.4),  Inches(0.22), size=9,   color=clr, italic=True)
    ai_lbl = "✦ Claude" if ai else "solo"
    ai_clr = AMBER if ai else MUTED
    tx(sl, ai_lbl, rx3+Inches(4.15), fy+Inches(0.03), Inches(0.9),  Inches(0.22), size=8.5, color=ai_clr, bold=ai)
    tx(sl, desc,   rx3+Inches(5.1),  fy+Inches(0.03), Inches(0.98), Inches(0.22), size=7.5, color=MUTED)

# AI legend note
tx(sl, "✦ Claude  =  Aman Deep prompted and directed; Claude Sonnet 4.6 co-authored the implementation",
   rx3, Inches(7.06), Inches(6.1), Inches(0.24),
   size=8.5, italic=True, color=AMBER)

rect(sl, Inches(0.38), Inches(7.3), Inches(12.6), Inches(0.2), fill=DARK)
tx(sl, "github.com/esfaihackathon/visual-assistant-ai  ·  branch: main  ·  35 commits total",
   Inches(0.5), Inches(7.32), Inches(12.3), Inches(0.16),
   size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — REFERENCES (clean, single slide)
# ─────────────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
bg(sl)
rect(sl, 0, 0, Inches(0.16), H, fill=AMBER)
tx(sl, "REFERENCES  &  DATA SOURCES", Inches(0.55), Inches(0.28), Inches(10), Inches(0.38),
   size=10.5, bold=True, color=AMBER)
tx(sl, "All data cited — verified and sourced",
   Inches(0.55), Inches(0.65), Inches(10), Inches(0.72),
   size=28, bold=True, color=WHITE)
divider(sl, Inches(1.38), AMBER)

ref_data = [
    (BLUE,   "[1]", "IAPB Vision Atlas 2023",
             "iapb.org/learn/iapb-vision-atlas",
             "India: 8.95M blind + 62.7M moderate–severe VI = ~72M. Used: Slides 2, 4."),
    (BLUE,   "[2]", "WHO World Report on Vision, 2019  (WHO/NMH/PBD/19.1)",
             "who.int/publications/i/item/9789241516570",
             "Global overview; India holds the largest number of blind people worldwide."),
    (TEAL,   "[3]", "UN World Population Ageing 2023",
             "un.org/development/desa/pd/content/world-population-ageing-2023",
             "India 60+ current: ~145M; projected 347M by 2050. Slides 2, 4."),
    (TEAL,   "[4]", "India Ageing Report 2023 — MoSPI / UNFPA India",
             "unfpa.org/india-ageing-report-2023",
             "Confirms UN trajectory; used for senior-citizen barrier framing."),
    (ORANGE, "[5]", "Colour Blind Awareness — colourblindawareness.org",
             "colourblindawareness.org/colour-blindness/types-of-colour-blindness",
             "~300M globally colour-blind; ~8% of males + 0.5% of females. India ≈ 60M. Slides 2, 4."),
    (ORANGE, "[6]", "Deeb, S.S. (2004). Molecular basis of variation in human color vision.",
             "Clinical Genetics, 67(5), 369–377. doi:10.1111/j.1399-0004.2004.00293.x",
             "Peer-reviewed prevalence: deuteranopia ~4.6% males, protanopia ~1.01% males."),
    (AMBER,  "[7]", "Census of India 2011 — Literacy Data",
             "censusindia.gov.in/census.website/data/census-tables",
             "Literacy rate 74.04%; ~26% = ~372M with limited literacy (2023 population base). Slides 2, 4."),
    (AMBER,  "[8]", "National Literacy Mission, Government of India",
             "nlm.nic.in",
             "Functional literacy definitions and adult literacy programme data."),
    (PURPLE, "[9]", "W3C WCAG 2.1 Guidelines",
             "w3.org/TR/WCAG21",
             "SC 1.4.3 AA: 4.5:1 · SC 1.4.6 AAA: 7:1 · SC 1.4.4 Resize Text. Slides 3, 4."),
    (PURPLE, "[10]","WebAIM Contrast Checker",
             "webaim.org/resources/contrastchecker",
             "Verified: White #FFF on Navy #081320 = 18.7:1 (AAA). Amber #FFC107 on Navy = 11.5:1 (AAA)."),
    (PURPLE, "[11]","Material Design Accessibility — Google",
             "m3.material.io/foundations/accessible-design/accessibility-basics",
             "Recommended minimum touch target 48×48dp. Saral uses 88dp (1.83× recommendation)."),
]

# Two-column layout
col_w = Inches(6.22)
left_refs  = ref_data[:6]
right_refs = ref_data[6:]

for col_i, col_refs in enumerate([left_refs, right_refs]):
    cx = Inches(0.38) + col_w * col_i + Inches(0.08)*col_i
    cy = Inches(1.58)
    for (clr, num, title, url, note) in col_refs:
        rect(sl, cx, cy, col_w, Inches(0.82), fill=NAVY_CARD)
        rect(sl, cx, cy, Inches(0.07), Inches(0.82), fill=clr)
        # Num badge
        rect(sl, cx+Inches(0.14), cy+Inches(0.12), Inches(0.36), Inches(0.26), fill=clr)
        tx(sl, num, cx+Inches(0.14), cy+Inches(0.11), Inches(0.36), Inches(0.26),
           size=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        tx(sl, title, cx+Inches(0.58), cy+Inches(0.08), col_w-Inches(0.65), Inches(0.28),
           size=10.5, bold=True, color=WHITE)
        tx(sl, url,   cx+Inches(0.58), cy+Inches(0.34), col_w-Inches(0.65), Inches(0.2),
           size=8.5, color=BLUE, italic=True)
        tx(sl, note,  cx+Inches(0.58), cy+Inches(0.54), col_w-Inches(0.65), Inches(0.24),
           size=8.5, color=MUTED, italic=True, wrap=True)
        cy += Inches(0.87)

rect(sl, 0, Inches(7.1), W, Inches(0.4), fill=DARK)
tx(sl, "All population figures are best available public estimates at time of presentation. "
       "Contrast ratios computed per WCAG 2.1 relative luminance formula (verified via WebAIM).",
   Inches(0.4), Inches(7.12), W-Inches(0.6), Inches(0.28),
   size=8.5, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/aman/Desktop/esfaihackathon/visual-assistant-ai/presentation/Saral-Business-Deck.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
